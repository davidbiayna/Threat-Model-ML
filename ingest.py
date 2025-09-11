# src/ingest.py
import os, re, glob
from typing import List
import pandas as pd
from pptx import Presentation
from docx import Document

# ---------------- Config ----------------
STANDARD_COLS = [
    "ThreatID","System","ThreatType","Severity","Date",
    "Description","SourceFile","SourceThreatID"
]

SYN = {
    # IDs from source -> capture but do NOT trust for output ID
    "threatid":"SourceThreatID","id":"SourceThreatID","ticket":"SourceThreatID",
    "ref":"SourceThreatID","reference":"SourceThreatID","threatref":"SourceThreatID",
    "threat ref":"SourceThreatID","threat ref.":"SourceThreatID","sourcethreatid":"SourceThreatID",

    # System
    "system":"System","asset":"System","component":"System","service":"System",
    "app":"System","application":"System","platform":"System","module":"System",

    # Threat type / category
    "threattype":"ThreatType","type":"ThreatType","category":"ThreatType",
    "stride":"ThreatType","kill chain":"ThreatType","killchain":"ThreatType",

    # Severity
    "severity":"Severity","risk":"Severity","rating":"Severity",
    "priority":"Severity","impact":"Severity",

    # Date
    "date":"Date","identified":"Date","detecteddate":"Date","created":"Date","logged":"Date",

    # Threat text sources
    "description":"Description","details":"Details","notes":"Notes",
    "summary":"Summary","observation":"Observation","threat":"Threat",
    "threat title":"ThreatTitle","threat description":"ThreatDescription",

    # Process/action helpers
    "step no.":"StepNo","step no":"StepNo","action ref.":"ActionRef","action ref":"ActionRef",

    # Source filename variants
    "sourcefile":"SourceFile","source file":"SourceFile"
}

SEV_MAP = {
    "critical":"High","high":"High","sev1":"High","p1":"High",
    "med":"Medium","medium":"Medium","moderate":"Medium","sev2":"Medium","p2":"Medium",
    "low":"Low","informational":"Low","info":"Low","sev3":"Low","p3":"Low",
}

THREATY_RE = re.compile(
    r"(malicious|threat|attack|unauthori[sz]|compromis|exploit|ransom|steal|exfiltrat|"
    r"leak|intercept|spoof|phish|token|password|credential|api\s*key|secret|privileg|"
    r"lateral|bypass|inject|xss|csrf|denial of service|dos)",
    re.IGNORECASE
)
SECRET_RE = re.compile(r"(?i)\b(password|pass|pwd|secret|api[_\- ]?key|apikey|token|bearer)\s*[:=]\s*\S+")

# --------------- Helpers ---------------
def norm_header(h: str) -> str:
    key = re.sub(r"[^a-z0-9]+","", str(h).strip().lower())
    return SYN.get(key) or str(h).strip().title()

def norm_severity(s):
    if pd.isna(s): return None
    t = str(s).strip().lower()
    return SEV_MAP.get(t, s if s in ["Low","Medium","High"] else t.title())

def clean_text(t: str) -> str:
    if pd.isna(t): return ""
    t = str(t).replace("\u00a0"," ").replace("\u2013","-").replace("\u2014","-")
    t = re.sub(r"\s+"," ", t)
    t = SECRET_RE.sub(r"\1=[REDACTED]", t)
    return t.strip()

def _pad_or_truncate(vals: List[str], n: int):
    return (vals + [""]*n)[:n]

def infer_system_from_file(source_file: str) -> str:
    if not source_file: return ""
    name = re.sub(r"(?i)^threat\s*model[-\s:_]*","", source_file)
    name = re.sub(r"\.(pptx?|docx?|csv|xlsx?)$","", name, flags=re.I)
    return name.strip(" -_")

def coalesce_threat_text(row: pd.Series, all_cols: List[str]) -> str:
    threat_like = [c for c in all_cols if re.match(r"(?i)^threat(?!.*(mitigation|ref))", c)]
    fixed = [c for c in ["ThreatTitle","Threat","ThreatDescription","Description","Summary","Observation","Notes","Details"]
             if c in all_cols]
    cols = list(dict.fromkeys(threat_like + fixed))
    parts = []
    for c in cols:
        v = row.get(c)
        if isinstance(v,str) and v.strip():
            parts.append(clean_text(v))
    # de-dup repeated sentences
    seen, uniq = set(), []
    for p in parts:
        if p not in seen:
            uniq.append(p); seen.add(p)
    return clean_text(" ".join(uniq)) if uniq else ""

def looks_like_step_or_action(row: pd.Series) -> bool:
    step = str(row.get("StepNo","")).strip()
    action = str(row.get("ActionRef","")).strip()
    return bool(re.fullmatch(r"[sS]\d{1,3}", step) or re.fullmatch(r"[aA]\d{1,3}", action))

def has_tcode(row: pd.Series, all_cols: List[str]) -> bool:
    tref_cols = [c for c in all_cols if re.match(r"(?i)^Threat Ref", c)]
    for c in (["SourceThreatID"] + tref_cols):
        v = row.get(c)
        if isinstance(v,str) and re.search(r"\bT\d{1,3}\b", v, flags=re.I):
            return True
    return False

def is_threat_row(row: pd.Series, all_cols: List[str]) -> bool:
    desc = row.get("Description","")
    if isinstance(desc,str) and desc.strip():
        if looks_like_step_or_action(row):
            return bool(THREATY_RE.search(desc))
        return True
    return has_tcode(row, all_cols)

def fill_if_blank(a: pd.Series, b: pd.Series) -> pd.Series:
    # keep a unless it's NA or empty/whitespace, then take b
    a_stripped = a.astype(str)
    blank = a_stripped.str.strip().eq("") | a.isna()
    res = a.copy()
    res[blank] = b[blank]
    return res

# ------------- Extractors -------------
def tables_from_pptx(path: str) -> pd.DataFrame:
    rows = []
    prs = Presentation(path)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, sh in enumerate(slide.shapes, start=1):
            if not getattr(sh, "has_table", False):
                continue
            tbl = sh.table
            if len(tbl.rows) < 2:
                print(f"[WARN] {os.path.basename(path)} slide {slide_idx} table {shape_idx}: <2 rows, skipped")
                continue
            headers = [norm_header(c.text) for c in tbl.rows[0].cells]
            for r_idx in range(1, len(tbl.rows)):
                r = tbl.rows[r_idx]
                vals = [clean_text(c.text) for c in r.cells]
                rows.append(dict(zip(headers, _pad_or_truncate(vals, len(headers)))))
    df = pd.DataFrame(rows)
    if not df.empty:
        df["SourceFile"] = os.path.basename(path)
    return df

def tables_from_docx(path: str) -> pd.DataFrame:
    rows = []
    doc = Document(path)
    for ti, t in enumerate(doc.tables, start=1):
        if len(t.rows) < 2:
            print(f"[WARN] {os.path.basename(path)} table {ti}: <2 rows, skipped")
            continue
        headers = [norm_header(c.text) for c in t.rows[0].cells]
        for r_idx in range(1, len(t.rows)):
            r = t.rows[r_idx]
            vals = [clean_text(c.text) for c in r.cells]
            rows.append(dict(zip(headers, _pad_or_truncate(vals, len(headers)))))
    df = pd.DataFrame(rows)
    if not df.empty:
        df["SourceFile"] = os.path.basename(path)
    return df

def read_csv_any(path: str) -> pd.DataFrame:
    try:
        d = pd.read_csv(path)
    except UnicodeDecodeError:
        d = pd.read_csv(path, encoding="latin-1")
    d = d.applymap(lambda x: clean_text(x) if isinstance(x, str) else x)
    # if an incoming CSV already has a source file column, keep it; else add
    if "SourceFile" not in d.columns and "Sourcefile" not in d.columns:
        d["SourceFile"] = os.path.basename(path)
    return d

# ------------- Main -------------------
def main():
    parts = []

    for p in glob.glob("data/raw/**/*.pptx", recursive=True):
        try:
            d = tables_from_pptx(p)
            if not d.empty: parts.append(d)
        except Exception as e:
            print(f"[ERROR] PPTX parse failed for {os.path.basename(p)}: {e}")

    for p in glob.glob("data/raw/**/*.docx", recursive=True):
        try:
            d = tables_from_docx(p)
            if not d.empty: parts.append(d)
        except Exception as e:
            print(f"[ERROR] DOCX parse failed for {os.path.basename(p)}: {e}")

    for p in glob.glob("data/raw/**/*.csv", recursive=True):
        try:
            parts.append(read_csv_any(p))
        except Exception as e:
            print(f"[ERROR] CSV read failed for {os.path.basename(p)}: {e}")

    assert parts, "No rows extracted. Put .pptx/.docx/.csv under data/raw/."

    # Combine and normalise headers
    df = pd.concat(parts, ignore_index=True, sort=False)
    df.columns = [norm_header(c) for c in df.columns]

    # ---- Coalesce duplicate/variant columns (very important for your case) ----
    # SourceFile ← Sourcefile (if either present)
    if "SourceFile" not in df.columns:
        df["SourceFile"] = ""
    if "Sourcefile" in df.columns:
        df["SourceFile"] = fill_if_blank(df["SourceFile"], df["Sourcefile"])
        df.drop(columns=["Sourcefile"], inplace=True)

    # SourceThreatID ← Sourcethreatid / Threat Ref.* / Ref (already mapped by SYN, but coalesce just in case)
    if "SourceThreatID" not in df.columns:
        df["SourceThreatID"] = ""
    for col in list(df.columns):
        if col.lower() in {"sourcethreatid"} or col.lower().startswith("threat ref"):
            df["SourceThreatID"] = fill_if_blank(df["SourceThreatID"], df[col])
            if col != "SourceThreatID":
                df.drop(columns=[col], inplace=True)

    # ---- Build consolidated Description (threat text) ----
    df["Description"] = df.apply(lambda r: coalesce_threat_text(r, list(df.columns)), axis=1)

    # ---- Severity / Date normalisation (leave blank if unknown) ----
    df["Severity"] = df["Severity"].map(norm_severity) if "Severity" in df.columns else None
    df["Date"] = pd.to_datetime(df["Date"], errors="coerce") if "Date" in df.columns else pd.NaT

    # ---- System inference (after SourceFile coalesced) ----
    if "System" not in df.columns:
        df["System"] = ""
    df["System"] = df["System"].fillna("").astype(str).str.strip()
    df["System"] = df.apply(
        lambda r: r["System"] if r["System"] else infer_system_from_file(r.get("SourceFile","")),
        axis=1
    )
    # Fallback: Asset Title as system if still blank
    if "Asset Title" in df.columns:
        df["System"] = df.apply(
            lambda r: r["System"] if r["System"] else (str(r["Asset Title"]).strip() if pd.notna(r["Asset Title"]) else ""),
            axis=1
        )

    # ---- Keep every threat row: text present OR a T-code; S/A rows kept only if threaty ----
    all_cols = list(df.columns)
    keep_mask = df.apply(lambda r: is_threat_row(r, all_cols), axis=1)
    before, kept = len(df), int(keep_mask.sum())
    df = df[keep_mask].copy().reset_index(drop=True)

    # ---- Assign fresh sequential IDs ----
    df.insert(0, "ThreatID", [f"THR-{i+1:06d}" for i in range(len(df))])

    # ---- Ensure standard columns exist and order them first ----
    for c in STANDARD_COLS:
        if c not in df.columns:
            df[c] = None
    df = df[STANDARD_COLS + [c for c in df.columns if c not in STANDARD_COLS]]

    # ---- Outputs ----
    os.makedirs("data/processed", exist_ok=True)
    df.to_csv("data/processed/threat_models.csv", index=False)
    try:
        df.to_parquet("data/processed/threat_models.parquet", index=False)
    except Exception as e:
        print(f"[WARN] Could not write parquet: {e}")

    # artifacts summary (engine-agnostic)
    os.makedirs("artifacts", exist_ok=True)
    summary = pd.DataFrame({
        "dtype": df.dtypes.astype(str),
        "non_null": df.notna().sum(),
        "nulls": df.isna().sum(),
        "unique": df.nunique(dropna=True),
    }).reset_index().rename(columns={"index":"column"})
    summary.to_csv("artifacts/dataset_summary.csv", index=False)

    # data dictionary
    os.makedirs("Docs", exist_ok=True)
    with open("Docs/Data_Dictionary.md","w") as f:
        f.write("# Data Dictionary (auto-generated)\n\n")
        f.write("| Field | Dtype | Example |\n|---|---|---|\n")
        for c in ["ThreatID","System","ThreatType","Severity","Date","Description","SourceFile","SourceThreatID"]:
            ex = df[c].dropna().astype(str).head(1).tolist()
            f.write(f"| {c} | {str(df[c].dtype)} | {ex[0][:80] if ex else ''} |\n")

    print("\nSaved:")
    print("  - data/processed/threat_models.csv")
    print("  - data/processed/threat_models.parquet (if supported)")
    print("  - artifacts/dataset_summary.csv")
    print("  - Docs/Data_Dictionary.md")
    print(f"\nExtracted rows: {before} | Kept threats: {kept} | Dropped non-threat scaffolding: {before-kept}")
    print("Columns:", list(df.columns))

if __name__ == "__main__":
    main()